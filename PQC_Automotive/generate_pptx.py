from pptx import Presentation
from pptx.util import Inches, Pt
import traceback
import sys

print('Creating Presentation_PQC_Automotive.pptx')
try:
    # Try to load the newly supplied Actia template
    prs = Presentation('TEMPLATE-ACTIA 2024.pptx')
    print('Loaded TEMPLATE-ACTIA 2024.pptx.')
except Exception as e:
    print('Failed to load TEMPLATE-ACTIA 2024.pptx. Falling back to default.')
    print(traceback.format_exc())
    prs = Presentation()

# Helper to add a Title Slide
def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        for p in slide.placeholders:
            if not p.is_placeholder: continue
            if 'subtitle' in p.name.lower() or 'sous-titre' in p.name.lower():
                p.text = subtitle
                break
        else:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle

# Helper to add a Content Slide
def add_content_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    if slide.placeholders and len(slide.placeholders) > 1:
        body_placeholder = None
        for p in slide.placeholders:
             if p == slide.shapes.title: continue
             if p.is_placeholder:
                 body_placeholder = p
                 break
                 
        if body_placeholder:
            tf = body_placeholder.text_frame
            tf.clear()
            for i, point in enumerate(bullet_points):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = point
                p.level = 0
                if len(point) > 90:
                   p.font.size = Pt(16)
                else:
                   p.font.size = Pt(20)

# Slide 1: Title
add_title_slide('Securing Automotive Networks against the Quantum Threat', 'Post-Quantum Cryptography in UDS 0x29 Architecture')

# Slide 2: The Quantum Threat
add_content_slide('The Quantum Threat', [
    'Current automotive security relies heavily on RSA and ECC (Elliptic Curve Cryptography).',
    'In 1994, Shor proven mathematically that a large Quantum Computer breaks these algorithms instantly.',
    'When Cryptographically Relevant Quantum Computers (CRQCs) arrive, ECC and RSA offer zero security.',
    'Given 10-15 year vehicle lifespans, ECUs designed today are already at risk.'
])

# Slide 3: Why Not QKD?
add_content_slide('QKD vs PQC in Automotive', [
    'Quantum Key Distribution (QKD) secures keys using physical quantum mechanics (photons).',
    'Why it fails for automotive:',
    '  - Requires pristine fiber-optic cables, dedicated lasers, and expensive hardware.',
    '  - Impossible to run over copper CAN buses.',
    '  - Cannot provide digital signatures for OTA firmware updates.',
    'The mandated NIST Solution: Post-Quantum Cryptography (PQC).'
])

# Slide 4: Post-Quantum Cryptography (PQC)
add_content_slide('The PQC Solution', [
    'PQC relies on complex new mathematics (like Lattice-based crypto) resistant to quantum algorithms.',
    'No hardware needed — runs entirely on standard classic microcontrollers (e.g., STM32).',
    'NIST Standardized Algorithms:',
    '  - ML-DSA (Dilithium) for Digital Signatures (Authentication & OTA).',
    '  - ML-KEM (Kyber) for Key Encapsulation (Secure Sessions).',
    'The Engineering Challenge: PQC keys are 10x to 50x larger than classical ECC keys.'
])

# Slide 5: Project Phases
add_content_slide('Project Overview & Execution Phases', [
    'Phase 1: Automotive Transport Baseline (Overcome 8-byte CAN limits for 15KB data).',
    'Phase 2: Crypto Abstraction Layer (CAL) (Decouple UDS logic from cryptography).',
    'Phase 3: Classical UDS 0x29 State Machine (Establish timing and memory baselines).',
    'Phase 4: Post-Quantum Integration (Swap Classical for ML-DSA and ML-KEM).',
    'Phase 5: Hybrid Mode & GUI Demonstration (TouchGFX visual hot-swapping).'
])

# Slide 6: Current Advancements
add_content_slide('Current Advancements', [
    'Transport Layer (Done): ISO-TP extended framing successfully moves 15.3 KB of data over FDCAN.',
    'Crypto Abstraction Layer (Done): Vtable architecture wrapping mbedTLS (X.509, ECDH, AES-GCM).',
    'UDS 0x29 State Machine (Done):',
    '  - Zero-Copy memory overlays implemented (saving ~4KB RAM).',
    '  - NRC 0x78 (Response Pending) flow natively prevents P2 timeouts during 122ms crypto bounds.'
])

# Slide 7: Next Steps
add_content_slide('Immediate Next Steps', [
    'Physical Hardware Deployment:',
    '  - Flash the completed Transport, CAL, and UDS binaries onto the twin STM32H7 evaluation boards.',
    '  - Test the Classical UDS 0x29 Authentication flow over live FDCAN.',
    '  - Profile runtime execution speed and memory consumption.',
    '  - Validate that NRC 0x78 pending timeouts successfully hold the diagnostic channel open.',
    'Once classical baseline is validated, integrate PQClean (ML-DSA) into the CAL Vtable.'
])

try:
   prs.save('Presentation_PQC_Automotive.pptx')
   print('Successfully generated Presentation_PQC_Automotive.pptx')
except Exception as e:
   print('Failed to save presentation!')
   print(e)
   sys.exit(1)
